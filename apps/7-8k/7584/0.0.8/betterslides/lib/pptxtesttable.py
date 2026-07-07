from pptx import Presentation
from pptx.util import Inches

# Create a PowerPoint presentation
presentation = Presentation()

# Add a slide
slide = presentation.slides.add_slide(presentation.slide_layouts[1])  # Index 1 corresponds to a blank slide layout

# Define table dimensions
rows = 3
cols = 4

# Set table position and dimensions
left = Inches(1)
top = Inches(1)
width = Inches(8)
height = Inches(4)

# Add a table to the slide
table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Set table column widths
column_widths = [Inches(2), Inches(2), Inches(2), Inches(2)]
for i, width in enumerate(column_widths):
    table.columns[i].width = width

# Populate the table with data
data = [
    ['Header 1', 'Header 2', 'Header 3', 'Header 4'],
    ['Value 1', 'Value 2', 'Value 3', 'Value 4'],
    ['Value 5', 'Value 6', 'Value 7', 'Value 8']
]

for row in range(rows):
    for col in range(cols):
        table.cell(row, col).text = data[row][col]

# Save the presentation
presentation.save('output.pptx')
