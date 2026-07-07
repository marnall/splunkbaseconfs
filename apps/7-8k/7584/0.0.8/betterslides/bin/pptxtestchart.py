import os,sys 
#load own libs from ../lib
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', 'lib'))

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

# Create a PowerPoint presentation
presentation = Presentation()

# Add a slide
slide = presentation.slides.add_slide(presentation.slide_layouts[1])  # Index 1 corresponds to a blank slide layout

# Define chart data
categories = ['Category 1', 'Category 2', 'Category 3']
data = [40, 30, 20]

# Set chart position and dimensions
left = Inches(1)
top = Inches(1)
width = Inches(6)
height = Inches(4)

# Add chart data
chart_data = CategoryChartData()
chart_data.categories = categories
chart_data.add_series('Series 1', data)

# Add a chart to the slide
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.PIE, left, top, width, height, chart_data
).chart

# Set chart title
chart.has_title = True
chart.chart_title.text_frame.text = "Pie Chart Example"

# Plot the chart data
chart.plots[0].has_data_labels = True
data_labels = chart.plots[0].data_labels
data_labels.number_format = '0%'
#data_labels.position = 'outside_end'

abspath = os.path.abspath(__file__)
dname = os.path.dirname(abspath)
os.chdir(dname)
# Save the presentation
presentation.save('pptxtestchart.pptx')
